"""
    Author:         CaptCorpMURICA
    Project:        100DaysPython
    File:           module3_day34_powerpoint.py
    Creation Date:  6/28/2019, 3:35 PM
    Description:    Python Automation Program 6: PowerPoint Manipulations
"""

# The `python-pptx` package can be used to interact with PowerPoint presentations. When a presentation is read, a
# presentation object is created.
from pptx import Presentation

pres = "TheDoctor.pptx"
doctors = Presentation(pres)
print(f"doctors: {doctors}")

print("\n", "=" * 50, "\n")

# The presentation object contains the `.slides` method. Each slide in the presentation is a slide object which can be
# called to pull information contained in the PowerPoint deck.
for i, slide in enumerate(doctors.slides):
    print(f"{i}: {slide}")

print("\n", "=" * 50, "\n")

# The slide contains specific identifiers and names based on the format of the slide. Additionally, the elements in the
# slide are in an OpenXML format.
title_slide, content_slide = doctors.slides[0], doctors.slides[1]
print(f"Slide ID: {title_slide.slide_id}\nSlide Element: {title_slide.element}\nSlide Layout: "
      f"{title_slide.slide_layout.name}")
print(f"Slide ID: {content_slide.slide_id}\nSlide Element: {content_slide.element}\nSlide Layout: "
      f"{content_slide.slide_layout.name}")

print("\n", "=" * 50, "\n")

# Within each slide, there are objects that are used to contain content.
for i, slide in enumerate(doctors.slides, start=1):
    print(f"Slide {i}, Type: {slide.slide_layout.name}")
    for shape in slide.shapes:
        print(f"Shape Type: {shape.placeholder_format.idx}\tPlaceholder Name: {shape.name}")

print("\n", "=" * 50, "\n")

# The text within the placeholder can then be extracted and stored as a dictionary with the tile of the slide as the key
# and the values of the Content Placeholder as the key/value pairs in the nested dictionary. The resulting dictionary
# resembles a JSON file that contains the data from all of the relevant slides.
doctor = dict()
for i, slide in enumerate(doctors.slides, start=1):
    if slide.slide_layout.name == "Two Content":
        title = slide.shapes.title.text
        deets_dict = dict()
        for shape in slide.shapes:
            if shape.placeholder_format.idx == 1:
                details = shape.text.split("\n")
                for item in details:
                    deets_dict[item.split(": ")[0]] = item.split(": ")[1]
                doctor[title] = deets_dict
    else:
        continue

print(doctor)

print("\n", "=" * 50, "\n")

# The `python-pptx` package can also create presentations. There are functions for adding shapes, images, and formatting
# the presentation for a consistent output. As a simple example, the dictionary created by reading the initial deck can
# be used to automatically create a new deck. Just with other automation tasks, the initial development takes time to
# effectively create the desired output. This package can be used to automate report generation based off data in an
# Excel spreadsheet. While this package is still in beta, it is already powerful enough to automate many of these
# clerical tasks.
from datetime import datetime

# A blank presentation object is created.
new_pres = Presentation()

# The title slide is added which includes a creation date in the subtitle.
slide = new_pres.slides.add_slide(new_pres.slide_layouts[0])
slide.shapes.title.text = "This was created with python".title()
slide.placeholders[1].text = f"Created on {datetime.date(datetime.now())}".title()

# The contents of the dictionary are iterated over to create a new slide for each new doctor.
for i, title in enumerate(doctor):
    new_slide = new_pres.slides.add_slide(new_pres.slide_layouts[3])
    shapes = new_slide.shapes
    title_shape = shapes.title
    title_shape.text = title
    # The details of the Doctor are added to the left placeholder. The actor is added first and then the remaining
    # details are iterated over to add as a new paragraph.
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = f"{list(doctor[title].keys())[0]}: {doctor[title][list(doctor[title].keys())[0]]}"
    for j in range(1, len(doctor[title])):
        tf.add_paragraph().text = f"{list(doctor[title].keys())[j]}: {doctor[title][list(doctor[title].keys())[j]]}"

# After the content is added, the presentation is saved.
new_pres.save('auto_pres.pptx')
